#!/usr/bin/env python3
"""
Generate PPT for Smart NPC Solution — standalone, no template dependency.
Clean AWS-style design with custom styling.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, copy

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "smart-npc-solution.pptx")

# Slide dimensions: 13.33" x 7.5" (widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Color palette (AWS brand) ──
C_DARK    = RGBColor(0x23, 0x2F, 0x3E)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_ORANGE  = RGBColor(0xFF, 0x99, 0x00)
C_SQUID   = RGBColor(0xEC, 0x7E, 0x11)  # darker orange
C_GRAY_L  = RGBColor(0xF2, 0xF3, 0xF3)
C_GRAY_D  = RGBColor(0x54, 0x58, 0x5D)
C_BLUE    = RGBColor(0x00, 0x73, 0xBB)
C_BLUE_L  = RGBColor(0xE1, 0xF0, 0xFA)
C_GREEN   = RGBColor(0x1B, 0x66, 0x0F)
C_GREEN_L = RGBColor(0xE8, 0xF5, 0xE9)
C_RED     = RGBColor(0xCC, 0x00, 0x00)
C_RED_L   = RGBColor(0xFF, 0xEB, 0xEE)
C_PURPLE  = RGBColor(0x8C, 0x4F, 0xFF)
C_AMBER_L = RGBColor(0xFD, 0xF0, 0xD5)

FONT = "Calibri"      # universally available fallback
FONT_B = "Calibri"

# ── Helpers ──

def _tb(slide, l, t, w, h, text, sz=14, color=C_DARK, bold=False,
        align=PP_ALIGN.LEFT, name=FONT, wrap=True, valign=None):
    """Add text box. Returns the shape."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return box


def _multi(slide, l, t, w, h, lines, sz=14, color=C_DARK, spacing=Pt(4),
           name=FONT, bold=False, align=PP_ALIGN.LEFT, bullet=""):
    """Multi-line text box. Each element of `lines` becomes a paragraph."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}{line}" if bullet else line
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = name
        p.font.bold = bold
        p.alignment = align
        p.space_after = spacing
    return box


def _rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(1)):
    """Rectangle shape."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s


def _rrect(slide, l, t, w, h, fill=None, border=None, bw=Pt(1.5)):
    """Rounded rectangle."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s


def _oval(slide, l, t, w, h, fill=None, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def _line_h(slide, l, t, w, color=C_ORANGE, thickness=Pt(3)):
    """Horizontal line (thin rect)."""
    return _rect(slide, l, t, w, thickness, fill=color)


def _card(slide, l, t, w, h, title, body_lines, title_color=C_BLUE,
          bg=C_WHITE, border=None, title_sz=15, body_sz=12):
    """A card with title + bullet body."""
    bdr = border or title_color
    _rrect(slide, l, t, w, h, fill=bg, border=bdr)
    _tb(slide, l + Inches(0.15), t + Inches(0.12), w - Inches(0.3), Inches(0.35),
        title, sz=title_sz, color=title_color, bold=True, align=PP_ALIGN.CENTER)
    _multi(slide, l + Inches(0.15), t + Inches(0.55), w - Inches(0.3),
           h - Inches(0.65), body_lines, sz=body_sz, color=C_GRAY_D,
           align=PP_ALIGN.CENTER, spacing=Pt(3))


def new_slide(prs):
    """Add a blank slide."""
    layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(layout)


def _page_number(slide, num):
    """Small page number bottom-right."""
    _tb(slide, SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.45),
        Inches(0.6), Inches(0.3), str(num), sz=10, color=C_GRAY_D,
        align=PP_ALIGN.RIGHT)


def _section_title_bar(slide, text, sub=""):
    """Standard slide title area with orange underline."""
    _tb(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.55),
        text, sz=28, color=C_DARK, bold=True, name=FONT_B)
    _line_h(slide, Inches(0.8), Inches(0.88), Inches(1.5))
    if sub:
        _tb(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.35),
            sub, sz=13, color=C_GRAY_D)


# ═══════════════════════════════════════════════════════════════════════════
# Slide builders
# ═══════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Cover page."""
    slide = new_slide(prs)
    # Dark background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_DARK

    # Orange accent bar top
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill=C_ORANGE)

    # Title
    _tb(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(1.0),
        "游戏智能 NPC 解决方案", sz=42, color=C_WHITE, bold=True,
        name=FONT_B, align=PP_ALIGN.LEFT)
    # Subtitle
    _tb(slide, Inches(1.2), Inches(3.2), Inches(10), Inches(0.6),
        "基于 Amazon Bedrock & AgentCore 的 AI 驱动动态任务系统",
        sz=20, color=C_ORANGE, name=FONT)
    # Thin orange line
    _line_h(slide, Inches(1.2), Inches(4.0), Inches(3), C_ORANGE, Pt(2))
    # Bottom info
    _tb(slide, Inches(1.2), Inches(5.5), Inches(6), Inches(0.4),
        "Amazon Web Services  |  解决方案详细介绍", sz=14, color=C_GRAY_D)
    _page_number(slide, 1)


def slide_agenda(prs):
    slide = new_slide(prs)
    _section_title_bar(slide, "Agenda")

    items = [
        ("01", "客户痛点与应用场景", "传统 NPC 的局限与 AI 智能 NPC 的价值"),
        ("02", "AWS 核心产品介绍", "Amazon Bedrock & AgentCore"),
        ("03", "方案技术架构", "整体架构、AI 模块架构、行为日志时序"),
        ("04", "Demo 演示", "智能 NPC 实机演示视频"),
        ("05", "价格分析", "AgentCore & Bedrock 按月成本估算"),
        ("06", "方案总结", "方案优势对比与正式生产架构"),
    ]
    y0 = Inches(1.5)
    for i, (num, title, desc) in enumerate(items):
        y = y0 + Inches(i * 0.9)
        # number circle
        o = _oval(slide, Inches(1.0), y + Inches(0.02), Inches(0.48), Inches(0.48),
                  fill=C_ORANGE)
        tf = o.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = num; p.font.size = Pt(14); p.font.color.rgb = C_WHITE
        p.font.bold = True; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # title + desc
        _tb(slide, Inches(1.7), y, Inches(8), Inches(0.35),
            title, sz=17, color=C_DARK, bold=True)
        _tb(slide, Inches(1.7), y + Inches(0.35), Inches(8), Inches(0.3),
            desc, sz=12, color=C_GRAY_D)

    _page_number(slide, 2)


def slide_section(prs, num, title, subtitle, page):
    """Section divider: dark bg, large number."""
    slide = new_slide(prs)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C_DARK
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill=C_ORANGE)

    _tb(slide, Inches(1), Inches(2.2), Inches(2), Inches(1.2),
        num, sz=60, color=C_ORANGE, bold=True, name=FONT_B, align=PP_ALIGN.LEFT)
    _tb(slide, Inches(1), Inches(3.5), Inches(10), Inches(0.7),
        title, sz=34, color=C_WHITE, bold=True, name=FONT_B)
    _tb(slide, Inches(1), Inches(4.3), Inches(10), Inches(0.5),
        subtitle, sz=16, color=C_GRAY_D)
    _line_h(slide, Inches(1), Inches(5.0), Inches(2.5), C_ORANGE, Pt(2))
    _page_number(slide, page)


# ── Section 1: Pain points ──

def slide_pain_points(prs):
    slide = new_slide(prs)
    _section_title_bar(slide, "客户痛点")

    # Left: Traditional
    _rrect(slide, Inches(0.6), Inches(1.3), Inches(5.6), Inches(5.3),
           fill=C_RED_L, border=C_RED)
    _tb(slide, Inches(0.8), Inches(1.4), Inches(5.2), Inches(0.4),
        "传统 NPC 实现方式", sz=18, color=C_RED, bold=True)
    _line_h(slide, Inches(0.8), Inches(1.85), Inches(1.0), C_RED, Pt(2))

    trad = [
        "NPC 对话基于固定脚本，缺乏变化和沉浸感",
        "任务由策划团队手动录入，开发周期长",
        "玩家行为无法影响 NPC 反应，体验千篇一律",
        "内容更新需要发版，无法动态调整",
        "策划人力成本高，难以规模化生成内容",
    ]
    _multi(slide, Inches(0.9), Inches(2.1), Inches(5.1), Inches(4.0),
           trad, sz=14, color=C_GRAY_D, bullet="\u2716  ", spacing=Pt(12))

    # Arrow
    _tb(slide, Inches(5.9), Inches(3.5), Inches(1.0), Inches(0.5),
        "\u276f\u276f", sz=30, color=C_ORANGE, bold=True, align=PP_ALIGN.CENTER)

    # Right: AI
    _rrect(slide, Inches(6.8), Inches(1.3), Inches(5.6), Inches(5.3),
           fill=C_GREEN_L, border=C_GREEN)
    _tb(slide, Inches(7.0), Inches(1.4), Inches(5.2), Inches(0.4),
        "AI 智能 NPC 解决方案", sz=18, color=C_GREEN, bold=True)
    _line_h(slide, Inches(7.0), Inches(1.85), Inches(1.0), C_GREEN, Pt(2))

    ai = [
        "NPC 对话由大语言模型动态生成，自然且个性化",
        "任务根据玩家行为自动生成，无需策划手动设计",
        "基于玩家行为日志分析，NPC 动态调整任务方向",
        "实时更新，无需发版即可优化 NPC 行为",
        "显著降低策划人力成本，提高内容生成效率",
    ]
    _multi(slide, Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.0),
           ai, sz=14, color=C_GREEN, bullet="\u2714  ", spacing=Pt(12))

    _page_number(slide, 4)


def slide_scenario_flow(prs):
    """Core scenario flow from design.md."""
    slide = new_slide(prs)
    _section_title_bar(slide, "应用场景 —— 核心场景流程",
                       "玩家在场景中与怪物战斗、与 NPC 对话并接取 AI 动态生成的任务")

    # Build a flow using simple shapes
    def box(l, t, w, h, text, fill, fc=C_WHITE, sz=13):
        _rrect(slide, l, t, w, h, fill=fill)
        _tb(slide, l, t, w, h, text, sz=sz, color=fc, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    def diamond(l, t, w, h, text, fill, fc=C_WHITE, sz=12):
        s = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = fill
        s.line.fill.background()
        _tb(slide, l + Inches(0.1), t + Inches(0.15), w - Inches(0.2),
            h - Inches(0.3), text, sz=sz, color=fc, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    def arrow_text(l, t, text, sz=10):
        _tb(slide, l, t, Inches(1.5), Inches(0.2), text, sz=sz, color=C_GRAY_D,
            align=PP_ALIGN.CENTER)

    # Row 0: Start
    box(Inches(5.2), Inches(1.5), Inches(2.6), Inches(0.55),
        "玩家进入场景", C_DARK)
    arrow_text(Inches(5.7), Inches(2.08), "\u2193")

    # Row 1: Move
    box(Inches(5.2), Inches(2.25), Inches(2.6), Inches(0.55),
        "键盘操控角色移动", C_BLUE)
    arrow_text(Inches(5.7), Inches(2.83), "\u2193")

    # Row 2: Decision
    diamond(Inches(5.3), Inches(3.0), Inches(2.4), Inches(1.1),
            "遭遇对象?", C_GRAY_D)

    # Left branch label
    _tb(slide, Inches(3.5), Inches(3.2), Inches(1.5), Inches(0.25),
        "\u2190  怪物", sz=11, color=C_RED, bold=True, align=PP_ALIGN.RIGHT)
    # Right branch label
    _tb(slide, Inches(8.0), Inches(3.2), Inches(1.5), Inches(0.25),
        "NPC  \u2192", sz=11, color=C_BLUE, bold=True, align=PP_ALIGN.LEFT)

    # ── Left: Battle ──
    box(Inches(1.0), Inches(3.2), Inches(2.2), Inches(0.55),
        "自动进入战斗", C_RED)

    diamond(Inches(1.1), Inches(4.1), Inches(2.0), Inches(0.9),
            "战斗结果?", C_RED, sz=11)

    # Victory
    _tb(slide, Inches(0.2), Inches(4.35), Inches(0.9), Inches(0.2),
        "胜利 \u2193", sz=10, color=C_GREEN, bold=True, align=PP_ALIGN.CENTER)
    box(Inches(0.0), Inches(5.2), Inches(1.8), Inches(0.8),
        "战斗结算\n经验/金币/掉落", C_GREEN, sz=11)

    # Defeat
    _tb(slide, Inches(3.0), Inches(4.35), Inches(0.9), Inches(0.2),
        "失败 \u2193", sz=10, color=RGBColor(0x99,0x00,0x00), bold=True,
        align=PP_ALIGN.CENTER)
    box(Inches(2.1), Inches(5.2), Inches(1.8), Inches(0.8),
        "回到出生点\n恢复满血", RGBColor(0x99,0x00,0x00), sz=11)

    # ── Right: NPC ──
    box(Inches(9.5), Inches(3.2), Inches(2.5), Inches(0.55),
        "触发 NPC 对话", C_BLUE)

    # AI highlight
    _rrect(slide, Inches(9.2), Inches(4.1), Inches(3.2), Inches(1.2),
           fill=C_AMBER_L, border=C_ORANGE)
    _tb(slide, Inches(9.3), Inches(4.15), Inches(3.0), Inches(0.3),
        "\u2728 AI 核心能力", sz=14, color=C_SQUID, bold=True,
        align=PP_ALIGN.CENTER)
    _tb(slide, Inches(9.3), Inches(4.5), Inches(3.0), Inches(0.7),
        "NPC Agent 分析玩家行为日志\n动态生成个性化任务", sz=12,
        color=C_GRAY_D, align=PP_ALIGN.CENTER)

    box(Inches(9.5), Inches(5.6), Inches(2.5), Inches(0.55),
        "玩家接取任务", C_GREEN)

    # Bottom: task check
    diamond(Inches(5.3), Inches(6.0), Inches(2.4), Inches(1.0),
            "任务条件满足?", C_GRAY_D, sz=11)

    box(Inches(8.2), Inches(6.2), Inches(2.0), Inches(0.55),
        "完成 \u2192 领取奖励", C_GREEN, sz=11)

    _page_number(slide, 5)


# ── Section 2: AWS Products ──

def slide_bedrock(prs):
    slide = new_slide(prs)
    _section_title_bar(slide, "Amazon Bedrock",
                       "全托管的生成式 AI 服务，提供来自领先 AI 公司的高性能基础模型")

    cards = [
        ("多模型选择", ["提供 Claude、Llama 等", "多家顶级模型", "按需选择最佳模型"]),
        ("Converse API", ["统一对话接口", "简化多模型集成", "支持 Tool Use"]),
        ("安全合规", ["数据不用于模型训练", "VPC 私有端点", "SOC / ISO 合规"]),
        ("按量计费", ["按 Token 计费", "无最低消费", "成本可预测"]),
    ]
    for i, (t, b) in enumerate(cards):
        x = Inches(0.5 + i * 3.1)
        _card(slide, x, Inches(1.7), Inches(2.8), Inches(2.2),
              t, b, title_color=C_BLUE, bg=C_BLUE_L, border=C_BLUE)

    # Model info box
    _rrect(slide, Inches(0.5), Inches(4.3), Inches(12.0), Inches(2.7),
           fill=C_AMBER_L, border=C_ORANGE)
    _tb(slide, Inches(0.8), Inches(4.4), Inches(6), Inches(0.4),
        "本方案使用的模型", sz=18, color=C_SQUID, bold=True)
    _line_h(slide, Inches(0.8), Inches(4.85), Inches(1.0), C_ORANGE, Pt(2))

    _tb(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.35),
        "Claude 3.5 Haiku  (us.anthropic.claude-3-5-haiku-20241022-v1:0)",
        sz=16, color=C_DARK, bold=True)

    model_details = [
        "Anthropic 最新轻量级模型，推理速度极快、成本极低",
        "非常适合游戏场景中的实时 NPC 对话和任务生成",
        "单次调用延迟约 3.5 秒，满足游戏交互的实时性要求",
        "200K Token 上下文窗口，可注入丰富的游戏数据作为 context",
    ]
    _multi(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(1.3),
           model_details, sz=13, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(6))

    _page_number(slide, 7)


def slide_agentcore(prs):
    slide = new_slide(prs)
    _section_title_bar(slide, "Amazon Bedrock AgentCore",
                       "全托管的 AI Agent 部署与运行服务，支持自定义 Agent 容器化部署")

    cards = [
        ("容器化部署", ["支持自定义 Docker 镜像", "灵活部署任意框架", "ARM64 架构支持"]),
        ("VPC 私有模式", ["Agent 运行在客户 VPC", "数据不出网络边界", "私有子网隔离"]),
        ("弹性扩缩", ["按请求量自动扩缩", "无需管理基础设施", "高可用架构"]),
        ("Runtime Endpoint", ["标准化 API 端点", "简化调用集成", "内置监控日志"]),
    ]
    for i, (t, b) in enumerate(cards):
        x = Inches(0.5 + i * 3.1)
        _card(slide, x, Inches(1.7), Inches(2.8), Inches(2.2),
              t, b, title_color=C_BLUE, bg=C_BLUE_L, border=C_BLUE)

    _rrect(slide, Inches(0.5), Inches(4.3), Inches(12.0), Inches(2.7),
           fill=C_AMBER_L, border=C_ORANGE)
    _tb(slide, Inches(0.8), Inches(4.4), Inches(6), Inches(0.4),
        "本方案中的使用方式", sz=18, color=C_SQUID, bold=True)
    _line_h(slide, Inches(0.8), Inches(4.85), Inches(1.0), C_ORANGE, Pt(2))

    usage = [
        "NPC Agent 以 Docker 容器形式部署到 AgentCore Runtime (ARM64)",
        "Agent 运行在 VPC 私有子网，通过 NAT Gateway 访问 Bedrock 和 DynamoDB",
        "游戏服务器通过 InvokeAgentRuntime API 调用 NPC Agent",
        "Agent 预获取 6 张 DynamoDB 表数据，单次调用 Bedrock Converse API 完成推理",
    ]
    _multi(slide, Inches(0.8), Inches(5.1), Inches(11), Inches(1.5),
           usage, sz=13, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(6))

    _page_number(slide, 8)


# ── Section 3: Architecture ──

def slide_arch_highlevel(prs):
    slide = new_slide(prs)
    _section_title_bar(slide, "整体架构 (High Level)",
                       "详细架构图请参见 architecture.drawio — Page 1")

    # Developer
    _rrect(slide, Inches(0.3), Inches(3.3), Inches(1.5), Inches(0.7),
           fill=C_DARK)
    _tb(slide, Inches(0.3), Inches(3.3), Inches(1.5), Inches(0.7),
        "开发者终端", sz=12, color=C_WHITE, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Arrow SSM
    _tb(slide, Inches(1.8), Inches(3.35), Inches(1.2), Inches(0.3),
        "SSM 端口转发 \u2192", sz=10, color=C_GRAY_D, align=PP_ALIGN.CENTER)

    # VPC container
    _rrect(slide, Inches(3.0), Inches(1.5), Inches(4.5), Inches(5.0),
           border=C_BLUE, bw=Pt(2))
    _tb(slide, Inches(3.1), Inches(1.55), Inches(3), Inches(0.25),
        "AWS VPC  \u2014  私有子网", sz=11, color=C_BLUE, bold=True)

    # EC2
    _rrect(slide, Inches(3.2), Inches(2.0), Inches(4.1), Inches(4.2),
           fill=C_BLUE_L, border=C_BLUE, bw=Pt(1))
    _tb(slide, Inches(3.3), Inches(2.05), Inches(3), Inches(0.25),
        "Amazon EC2 (Node.js + Express)", sz=11, color=C_BLUE, bold=True)

    modules = [
        "Express.static  \u2014  前端静态资源",
        "REST API / WebSocket",
        "BattleSystem  \u2014  战斗系统",
        "TaskManager  \u2014  任务奖励结算",
        "EventEmitter  \u2014  事件发射器",
    ]
    _multi(slide, Inches(3.4), Inches(2.5), Inches(3.8), Inches(3.2),
           modules, sz=12, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(6))

    # AI Module
    _rrect(slide, Inches(8.0), Inches(1.5), Inches(4.5), Inches(3.0),
           fill=C_AMBER_L, border=C_ORANGE, bw=Pt(2))
    _tb(slide, Inches(8.1), Inches(1.55), Inches(3), Inches(0.25),
        "AI 模块", sz=11, color=C_SQUID, bold=True)

    ai = [
        "AgentCore Runtime (VPC 私有模式)",
        "NPC Agent (Docker 容器)",
        "Bedrock Converse API",
        "Claude 3.5 Haiku",
    ]
    _multi(slide, Inches(8.2), Inches(2.0), Inches(4.0), Inches(2.2),
           ai, sz=12, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(6))

    # Arrow EC2 -> AI
    _tb(slide, Inches(7.1), Inches(2.5), Inches(1.1), Inches(0.3),
        "NPC 对话请求 \u2192", sz=9, color=C_SQUID, bold=True, align=PP_ALIGN.CENTER)

    # DynamoDB
    _rrect(slide, Inches(8.0), Inches(4.8), Inches(4.5), Inches(1.7),
           fill=C_GREEN_L, border=C_GREEN, bw=Pt(2))
    _tb(slide, Inches(8.1), Inches(4.85), Inches(3), Inches(0.25),
        "Amazon DynamoDB", sz=11, color=C_GREEN, bold=True)
    ddb = ["Players  |  Monsters  |  NPCs  |  Items",
           "Tasks  |  PlayerEventSummary"]
    _multi(slide, Inches(8.2), Inches(5.3), Inches(4.0), Inches(1.0),
           ddb, sz=12, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(4))

    # Arrows
    _tb(slide, Inches(7.1), Inches(5.2), Inches(1.1), Inches(0.3),
        "\u2190 读写 \u2192", sz=9, color=C_GREEN, bold=True, align=PP_ALIGN.CENTER)

    _page_number(slide, 10)


def slide_arch_ai_module(prs):
    slide = new_slide(prs)
    _section_title_bar(slide, "AI 模块架构 — NPC Agent 内部流程",
                       "详细架构图请参见 architecture.drawio — Page 3")

    steps = [
        ("1. 数据预获取", "并行读取 6 张\nDynamoDB 表", C_BLUE),
        ("2. Prompt 组装", "精简数据 +\nSystem Prompt", C_GRAY_D),
        ("3. LLM 调用", "单次 Bedrock\nConverse API", C_ORANGE),
        ("4. JSON 解析", "解析 dialogue\n+ task", C_GRAY_D),
        ("5. 任务校验", "validate_task\n6 项校验规则", C_RED),
        ("6. 写入 DB", "校验通过后\n写入 Tasks 表", C_GREEN),
    ]
    for i, (t, d, c) in enumerate(steps):
        x = Inches(0.25 + i * 2.15)
        _rrect(slide, x, Inches(1.5), Inches(1.9), Inches(1.7),
               fill=C_WHITE, border=c, bw=Pt(2))
        # colored top bar
        _rect(slide, x, Inches(1.5), Inches(1.9), Inches(0.06), fill=c)
        _tb(slide, x + Inches(0.1), Inches(1.6), Inches(1.7), Inches(0.35),
            t, sz=13, color=c, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, x + Inches(0.1), Inches(2.05), Inches(1.7), Inches(0.9),
            d, sz=11, color=C_GRAY_D, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            _tb(slide, x + Inches(1.85), Inches(2.0), Inches(0.4), Inches(0.4),
                "\u2192", sz=20, color=C_ORANGE, bold=True, align=PP_ALIGN.CENTER)

    # Data sources
    _tb(slide, Inches(0.5), Inches(3.5), Inches(10), Inches(0.35),
        "预获取数据源 (6 张 DynamoDB 表)", sz=16, color=C_DARK, bold=True)
    _line_h(slide, Inches(0.5), Inches(3.88), Inches(1.0), C_GREEN, Pt(2))

    tables = [
        ("Players", "玩家状态"),
        ("PlayerEvent\nSummary", "最近 20 条\n行为日志"),
        ("Monsters", "怪物字典"),
        ("Items", "道具字典"),
        ("NPCs", "NPC 字典"),
        ("Tasks", "已有任务"),
    ]
    for i, (n, d) in enumerate(tables):
        x = Inches(0.25 + i * 2.15)
        _rrect(slide, x, Inches(4.05), Inches(1.9), Inches(1.0),
               fill=C_GREEN_L, border=C_GREEN)
        _tb(slide, x + Inches(0.05), Inches(4.1), Inches(1.8), Inches(0.4),
            n, sz=11, color=C_GREEN, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, x + Inches(0.05), Inches(4.5), Inches(1.8), Inches(0.4),
            d, sz=10, color=C_GRAY_D, align=PP_ALIGN.CENTER)

    # Validation rules
    _tb(slide, Inches(0.5), Inches(5.25), Inches(5), Inches(0.35),
        "任务校验规则 (validate_task)", sz=16, color=C_DARK, bold=True)
    _line_h(slide, Inches(0.5), Inches(5.63), Inches(1.0), C_RED, Pt(2))

    rules = [
        "JSON 结构完整性 (title / description / conditions / awards)",
        "条件类型合法性 (kill_monster / collect_item / talk_to_npc / use_item)",
        "target_id 在字典表中存在",
        "奖励 item_id 在字典表中存在",
        "数值范围校验 (金币 1-1000, 经验 1-500)",
        "任务去重 (不与进行中任务重复)",
    ]
    _multi(slide, Inches(0.5), Inches(5.75), Inches(5.5), Inches(1.6),
           rules, sz=11, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(2))

    # Performance box
    _rrect(slide, Inches(7.0), Inches(5.25), Inches(5.5), Inches(2.0),
           fill=C_AMBER_L, border=C_ORANGE)
    _tb(slide, Inches(7.2), Inches(5.3), Inches(3), Inches(0.3),
        "性能特征", sz=15, color=C_SQUID, bold=True)
    _line_h(slide, Inches(7.2), Inches(5.63), Inches(0.8), C_ORANGE, Pt(2))
    perf = [
        "数据预获取:  ~200ms",
        "LLM 推理 (Claude 3.5 Haiku):  ~3.5s",
        "校验 + 写入:  ~100ms",
        "总耗时:  ~4 秒",
    ]
    _multi(slide, Inches(7.2), Inches(5.75), Inches(5.0), Inches(1.3),
           perf, sz=12, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(4))

    _page_number(slide, 11)


def slide_sequence(prs):
    """Behavior log sequence diagram."""
    slide = new_slide(prs)
    _section_title_bar(slide, "行为日志采集与供给时序",
                       "Design 文档 5.2.1 节 — 行为日志驱动的动态任务生成")

    # Participant headers
    parts = [
        ("玩家", Inches(0.3), C_DARK),
        ("游戏服务", Inches(2.3), C_BLUE),
        ("Kinesis", Inches(4.5), C_ORANGE),
        ("S3", Inches(6.3), C_GREEN),
        ("DynamoDB", Inches(8.3), C_BLUE),
        ("NPC Agent", Inches(10.5), C_SQUID),
    ]
    for name, x, color in parts:
        _rrect(slide, x, Inches(1.4), Inches(1.6), Inches(0.5), fill=color)
        _tb(slide, x, Inches(1.4), Inches(1.6), Inches(0.5),
            name, sz=12, color=C_WHITE, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # lifeline
        _rect(slide, x + Inches(0.78), Inches(1.9), Pt(1), Inches(5.0),
              fill=RGBColor(0xBB, 0xBB, 0xBB))

    # Messages
    msgs = [
        (0.3, 2.3, "1. 玩家行为 (战斗/道具/...)", 2.2, C_DARK),
        (2.3, 4.5, "2. 发送完整事件 JSON", 2.8, C_ORANGE),
        (4.5, 6.3, "3. Firehose 持久化全量日志", 3.3, C_GREEN),
        (2.3, 8.3, "4. 写入最近事件摘要 (保留50条)", 3.8, C_BLUE),
        (0.3, 2.3, "5. 接近 NPC, 触发对话", 4.5, C_DARK),
        (2.3, 10.5, "6. InvokeAgentRuntime (player_id, npc_id)", 5.0, C_SQUID),
        (10.5, 8.3, "7. get_player_events(limit=20)", 5.5, C_BLUE),
        (8.3, 10.5, "8. 返回最近 20 条行为日志", 6.0, C_SQUID),
    ]
    for frm, to, text, y, color in msgs:
        mn = min(frm, to)
        mx = max(frm, to)
        # label
        _tb(slide, Inches(mn + 0.8), Inches(y - 0.18), Inches(mx - mn),
            Inches(0.22), text, sz=9, color=color, align=PP_ALIGN.CENTER)
        # line
        _rect(slide, Inches(mn + 0.8), Inches(y + 0.04),
              Inches(mx - mn), Pt(1.5), fill=color)

    _page_number(slide, 12)


# ── Section 4: Demo ──

def slide_demo(prs):
    slide = new_slide(prs)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C_DARK
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill=C_ORANGE)

    _tb(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(0.8),
        "Demo 演示", sz=42, color=C_WHITE, bold=True,
        name=FONT_B, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(0.5),
        "智能 NPC 实机演示视频", sz=20, color=C_ORANGE, align=PP_ALIGN.CENTER)

    _rrect(slide, Inches(3), Inches(4.5), Inches(7.3), Inches(0.6),
           border=C_GRAY_D)
    _tb(slide, Inches(3), Inches(4.5), Inches(7.3), Inches(0.6),
        "[ 请在此处插入 Demo 视频 ]", sz=14, color=C_GRAY_D,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    _page_number(slide, 14)


# ── Section 5: Pricing ──

def slide_pricing(prs):
    slide = new_slide(prs)
    _section_title_bar(slide, "价格分析 — 月度成本估算",
                       "以下估算基于 Demo 使用场景 (约 1000 次 NPC 对话/月)")

    # AgentCore card
    _rrect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(4.2),
           fill=C_WHITE, border=C_ORANGE, bw=Pt(2))
    _rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.55), fill=C_ORANGE)
    _tb(slide, Inches(0.6), Inches(1.52), Inches(5.5), Inches(0.5),
        "Amazon Bedrock AgentCore", sz=18, color=C_WHITE, bold=True)

    ac_lines = [
        "计费模型:  按 Agent 运行时间 (容器实例) 计费",
        "",
        "Runtime 实例 (arm64, 0.5 vCPU + 1GB):",
        "    ~$0.01 / 小时",
        "    月度 (24x7):  ~$7.20",
        "",
        "Endpoint 请求费:  可忽略不计",
        "",
        "AgentCore 月度估算:  ~$7 - $15",
    ]
    y = Inches(2.2)
    for line in ac_lines:
        if line == "":
            y += Inches(0.12)
            continue
        is_total = line.startswith("AgentCore 月度")
        c = C_SQUID if is_total else C_GRAY_D
        b = is_total
        _tb(slide, Inches(0.8), y, Inches(5.2), Inches(0.25),
            line, sz=12, color=c, bold=b)
        y += Inches(0.25)

    # Bedrock card
    _rrect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(4.2),
           fill=C_WHITE, border=C_BLUE, bw=Pt(2))
    _rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.55), fill=C_BLUE)
    _tb(slide, Inches(6.9), Inches(1.52), Inches(5.5), Inches(0.5),
        "Amazon Bedrock (Claude 3.5 Haiku)", sz=18, color=C_WHITE, bold=True)

    br_lines = [
        "计费模型:  按输入/输出 Token 计费",
        "",
        "Claude 3.5 Haiku 定价:",
        "    输入:  $0.80 / 百万 Token",
        "    输出:  $4.00 / 百万 Token",
        "",
        "Demo 场景 (1000 次对话/月):",
        "    月输入 ~200 万 Token =  $1.60",
        "    月输出 ~50 万 Token  =  $2.00",
        "",
        "Bedrock 月度估算:  ~$3.60",
    ]
    y = Inches(2.2)
    for line in br_lines:
        if line == "":
            y += Inches(0.12)
            continue
        is_total = line.startswith("Bedrock 月度")
        c = C_BLUE if is_total else C_GRAY_D
        b = is_total
        _tb(slide, Inches(7.1), y, Inches(5.2), Inches(0.25),
            line, sz=12, color=c, bold=b)
        y += Inches(0.25)

    # Total summary
    _rrect(slide, Inches(0.5), Inches(5.9), Inches(12.1), Inches(1.3),
           fill=C_AMBER_L, border=C_ORANGE, bw=Pt(2))
    _tb(slide, Inches(0.8), Inches(6.0), Inches(5), Inches(0.35),
        "月度总成本估算 (Demo 场景)", sz=17, color=C_SQUID, bold=True)
    _tb(slide, Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.3),
        "AgentCore: ~$7-15   |   Bedrock: ~$3.60   |   DynamoDB (按需): ~$1-5   |   其他 (EC2, VPC): 按现有基础设施",
        sz=12, color=C_GRAY_D)
    _tb(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.35),
        "AI 模块月度总成本 (AgentCore + Bedrock):  约 $10 - $20  ——  成本极低，适合 Demo 及中小规模生产",
        sz=14, color=C_SQUID, bold=True)

    _page_number(slide, 16)


# ── Section 6: Summary ──

def slide_comparison(prs):
    slide = new_slide(prs)
    _section_title_bar(slide, "方案总结 — 优势对比")

    # Table header
    cols = [
        (Inches(0.5), Inches(2.0), "对比维度"),
        (Inches(2.5), Inches(4.2), "传统方式"),
        (Inches(6.7), Inches(5.6), "AI 智能 NPC 方案"),
    ]
    for x, w, t in cols:
        _rect(slide, x, Inches(1.3), w, Inches(0.5), fill=C_DARK)
        _tb(slide, x, Inches(1.3), w, Inches(0.5),
            t, sz=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE)

    rows = [
        ("NPC 对话", "固定脚本，千篇一律", "LLM 动态生成，自然且个性化"),
        ("任务生成", "策划手动设计和录入", "AI 根据玩家行为自动生成"),
        ("内容更新", "需要发版更新", "实时动态调整，无需发版"),
        ("玩家体验", "缺乏沉浸感", "个性化体验，NPC 理解玩家处境"),
        ("开发成本", "大量策划人力，周期长", "一次开发，AI 持续生成内容"),
        ("可扩展性", "内容量与人力成正比", "AI 无限生成，边际成本趋零"),
        ("数据驱动", "无法利用玩家行为数据", "行为日志驱动决策，持续优化"),
        ("运营成本", "高 (人力为主)", "极低 (~$10-20/月 AI 服务费)"),
    ]
    for i, (dim, old, new) in enumerate(rows):
        y = Inches(1.8 + i * 0.55)
        bg = C_GRAY_L if i % 2 == 0 else C_WHITE

        for x, w in [(Inches(0.5), Inches(2.0)),
                      (Inches(2.5), Inches(4.2)),
                      (Inches(6.7), Inches(5.6))]:
            _rect(slide, x, y, w, Inches(0.52), fill=bg)

        _tb(slide, Inches(0.6), y + Inches(0.08), Inches(1.8), Inches(0.35),
            dim, sz=12, color=C_DARK, bold=True)
        _tb(slide, Inches(2.6), y + Inches(0.08), Inches(4.0), Inches(0.35),
            old, sz=11, color=C_RED)
        _tb(slide, Inches(6.8), y + Inches(0.08), Inches(5.4), Inches(0.35),
            new, sz=11, color=C_GREEN)

    _page_number(slide, 18)


def slide_prod_arch(prs):
    slide = new_slide(prs)
    _section_title_bar(slide, "正式生产架构（目标）",
                       "详细架构图请参见 architecture.drawio — Page 2")

    # Simplified production architecture
    # User -> CloudFront
    _rrect(slide, Inches(0.3), Inches(2.0), Inches(1.3), Inches(0.6),
           fill=C_DARK)
    _tb(slide, Inches(0.3), Inches(2.0), Inches(1.3), Inches(0.6),
        "用户", sz=11, color=C_WHITE, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    _tb(slide, Inches(1.6), Inches(2.05), Inches(0.5), Inches(0.3),
        "\u2192", sz=16, color=C_GRAY_D, align=PP_ALIGN.CENTER)

    _rrect(slide, Inches(2.0), Inches(2.0), Inches(1.5), Inches(0.6),
           fill=C_PURPLE)
    _tb(slide, Inches(2.0), Inches(2.0), Inches(1.5), Inches(0.6),
        "CloudFront", sz=11, color=C_WHITE, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # S3 Frontend
    _tb(slide, Inches(2.2), Inches(2.7), Inches(1.1), Inches(0.2),
        "\u2193 静态资源", sz=9, color=C_GRAY_D, align=PP_ALIGN.CENTER)
    _rrect(slide, Inches(2.0), Inches(2.95), Inches(1.5), Inches(0.5),
           fill=C_GREEN)
    _tb(slide, Inches(2.0), Inches(2.95), Inches(1.5), Inches(0.5),
        "S3 前端", sz=10, color=C_WHITE, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # API path
    _tb(slide, Inches(3.5), Inches(2.05), Inches(0.8), Inches(0.3),
        "/api /ws \u2192", sz=9, color=C_GRAY_D, align=PP_ALIGN.CENTER)

    # VPC
    _rrect(slide, Inches(4.2), Inches(1.5), Inches(4.3), Inches(4.5),
           border=C_BLUE, bw=Pt(2))
    _tb(slide, Inches(4.3), Inches(1.55), Inches(2), Inches(0.2),
        "VPC", sz=10, color=C_BLUE, bold=True)

    # ALB
    _rrect(slide, Inches(4.4), Inches(1.9), Inches(1.5), Inches(0.55),
           fill=C_PURPLE)
    _tb(slide, Inches(4.4), Inches(1.9), Inches(1.5), Inches(0.55),
        "ALB", sz=11, color=C_WHITE, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # EC2
    _rrect(slide, Inches(4.4), Inches(2.7), Inches(3.8), Inches(3.0),
           fill=C_BLUE_L, border=C_BLUE)
    _tb(slide, Inches(4.5), Inches(2.75), Inches(2), Inches(0.2),
        "EC2 (私有子网)", sz=10, color=C_BLUE, bold=True)
    _multi(slide, Inches(4.6), Inches(3.1), Inches(3.4), Inches(2.0),
           ["REST API / WebSocket", "战斗系统", "任务管理", "事件发射器"],
           sz=10, color=C_GRAY_D, bullet="\u2022 ", spacing=Pt(3))

    # AI Module
    _rrect(slide, Inches(9.0), Inches(1.5), Inches(3.5), Inches(2.3),
           fill=C_AMBER_L, border=C_ORANGE, bw=Pt(2))
    _tb(slide, Inches(9.1), Inches(1.55), Inches(2), Inches(0.2),
        "AI 模块", sz=10, color=C_SQUID, bold=True)
    _multi(slide, Inches(9.2), Inches(1.9), Inches(3.2), Inches(1.6),
           ["AgentCore Runtime", "NPC Agent (Docker)",
            "Bedrock Claude 3.5 Haiku"],
           sz=10, color=C_GRAY_D, bullet="\u2022 ", spacing=Pt(3))

    # Arrow EC2 -> AI
    _tb(slide, Inches(8.1), Inches(2.8), Inches(1.0), Inches(0.3),
        "NPC\u2192", sz=9, color=C_SQUID, align=PP_ALIGN.CENTER)

    # Data Pipeline
    _rrect(slide, Inches(9.0), Inches(4.1), Inches(3.5), Inches(1.5),
           fill=C_AMBER_L, border=C_ORANGE, bw=Pt(2))
    _tb(slide, Inches(9.1), Inches(4.15), Inches(2), Inches(0.2),
        "数据流管道", sz=10, color=C_SQUID, bold=True)
    _multi(slide, Inches(9.2), Inches(4.5), Inches(3.2), Inches(0.9),
           ["Kinesis Data Stream", "Kinesis Firehose \u2192 S3"],
           sz=10, color=C_GRAY_D, bullet="\u2022 ", spacing=Pt(3))

    # Arrow EC2 -> Pipeline
    _tb(slide, Inches(8.1), Inches(4.5), Inches(1.0), Inches(0.3),
        "事件\u2192", sz=9, color=C_SQUID, align=PP_ALIGN.CENTER)

    # DynamoDB
    _rrect(slide, Inches(4.2), Inches(6.2), Inches(3.5), Inches(0.7),
           fill=C_GREEN_L, border=C_GREEN)
    _tb(slide, Inches(4.2), Inches(6.2), Inches(3.5), Inches(0.7),
        "DynamoDB (6 tables)", sz=12, color=C_GREEN, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # S3 Logs
    _rrect(slide, Inches(9.0), Inches(6.2), Inches(3.5), Inches(0.7),
           fill=C_GREEN_L, border=C_GREEN)
    _tb(slide, Inches(9.0), Inches(6.2), Inches(3.5), Inches(0.7),
        "S3 行为日志归档", sz=12, color=C_GREEN, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Arrows down
    _tb(slide, Inches(5.5), Inches(5.8), Inches(0.5), Inches(0.3),
        "\u2193", sz=14, color=C_GREEN, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(10.3), Inches(5.7), Inches(0.5), Inches(0.3),
        "\u2193", sz=14, color=C_GREEN, align=PP_ALIGN.CENTER)

    _page_number(slide, 19)


def slide_thankyou(prs):
    slide = new_slide(prs)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C_DARK
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill=C_ORANGE)
    _rect(slide, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06),
          fill=C_ORANGE)

    _tb(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(1.0),
        "Thank You", sz=48, color=C_WHITE, bold=True,
        name=FONT_B, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(0), Inches(3.8), SLIDE_W, Inches(0.5),
        "Amazon Web Services", sz=18, color=C_ORANGE, align=PP_ALIGN.CENTER)
    _page_number(slide, 20)


# ═══════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Generating PPT (standalone, no template)...")

    slide_title(prs)       ; print("  [1]  Title")
    slide_agenda(prs)      ; print("  [2]  Agenda")
    slide_section(prs, "01", "客户痛点与应用场景",
                  "传统 NPC 的局限与 AI 智能 NPC 的价值", 3)
    print("  [3]  Section 1 divider")
    slide_pain_points(prs) ; print("  [4]  Pain points")
    slide_scenario_flow(prs); print("  [5]  Scenario flow")
    slide_section(prs, "02", "AWS 核心产品介绍",
                  "Amazon Bedrock & AgentCore", 6)
    print("  [6]  Section 2 divider")
    slide_bedrock(prs)     ; print("  [7]  Bedrock")
    slide_agentcore(prs)   ; print("  [8]  AgentCore")
    slide_section(prs, "03", "方案技术架构",
                  "整体架构、AI 模块架构、行为日志时序", 9)
    print("  [9]  Section 3 divider")
    slide_arch_highlevel(prs); print("  [10] High-level arch")
    slide_arch_ai_module(prs); print("  [11] AI module arch")
    slide_sequence(prs)    ; print("  [12] Sequence diagram")
    slide_section(prs, "04", "Demo 演示",
                  "智能 NPC 实机演示", 13)
    print("  [13] Section 4 divider")
    slide_demo(prs)        ; print("  [14] Demo placeholder")
    slide_section(prs, "05", "价格分析",
                  "AgentCore & Bedrock 按月成本估算", 15)
    print("  [15] Section 5 divider")
    slide_pricing(prs)     ; print("  [16] Pricing")
    slide_section(prs, "06", "方案总结",
                  "方案优势对比与正式生产架构", 17)
    print("  [17] Section 6 divider")
    slide_comparison(prs)  ; print("  [18] Comparison table")
    slide_prod_arch(prs)   ; print("  [19] Production arch")
    slide_thankyou(prs)    ; print("  [20] Thank You")

    prs.save(OUTPUT_PATH)
    print(f"\nDone! Saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
