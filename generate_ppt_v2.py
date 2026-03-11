#!/usr/bin/env python3
"""
Generate PPT using template.pptx as style source.
Deletes all 99 example slides, then creates new ones using the template's layouts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os, copy

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, "template.pptx")
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "smart-npc-solution-v2.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
C_DARK    = RGBColor(0x23, 0x2F, 0x3E)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_ORANGE  = RGBColor(0xFF, 0x99, 0x00)
C_SQUID   = RGBColor(0xEC, 0x7E, 0x11)
C_GRAY_L  = RGBColor(0xF2, 0xF3, 0xF3)
C_GRAY_D  = RGBColor(0x54, 0x58, 0x5D)
C_BLUE    = RGBColor(0x00, 0x73, 0xBB)
C_BLUE_L  = RGBColor(0xE1, 0xF0, 0xFA)
C_GREEN   = RGBColor(0x1B, 0x66, 0x0F)
C_GREEN_L = RGBColor(0xE8, 0xF5, 0xE9)
C_RED     = RGBColor(0xCC, 0x00, 0x00)
C_RED_L   = RGBColor(0xFF, 0xEB, 0xEE)
C_PURPLE  = RGBColor(0x8C, 0x4F, 0xFF)
C_AMBER_L = RGBColor(0xFD, 0xF0, 0xD5)

# ── Helpers ──

def get_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return None

def delete_all_slides(prs):
    """Remove all existing slides from the presentation."""
    while len(prs.slides._sldIdLst) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

def add_slide(prs, layout_name):
    layout = get_layout(prs, layout_name)
    if not layout:
        layout = get_layout(prs, 'Blank')
    return prs.slides.add_slide(layout)

def _tb(slide, l, t, w, h, text, sz=14, color=C_DARK, bold=False,
        align=PP_ALIGN.LEFT, wrap=True, valign=None):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = wrap
    if valign: tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = color
    p.font.bold = bold; p.alignment = align
    return box

def _multi(slide, l, t, w, h, lines, sz=14, color=C_DARK, spacing=Pt(4),
           bold=False, align=PP_ALIGN.LEFT, bullet=""):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}{line}" if bullet else line
        p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.bold = bold; p.alignment = align; p.space_after = spacing
    return box

def _rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if border: s.line.color.rgb = border; s.line.width = bw
    else: s.line.fill.background()
    return s

def _rrect(slide, l, t, w, h, fill=None, border=None, bw=Pt(1.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if border: s.line.color.rgb = border; s.line.width = bw
    else: s.line.fill.background()
    return s

def _oval(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    s.line.fill.background()
    return s

def _hline(slide, l, t, w, color=C_ORANGE, thickness=Pt(3)):
    return _rect(slide, l, t, w, thickness, fill=color)

def _card(slide, l, t, w, h, title, body_lines, title_color=C_BLUE,
          bg=C_WHITE, border=None, title_sz=15, body_sz=12):
    bdr = border or title_color
    _rrect(slide, l, t, w, h, fill=bg, border=bdr)
    _tb(slide, l+Inches(0.15), t+Inches(0.12), w-Inches(0.3), Inches(0.35),
        title, sz=title_sz, color=title_color, bold=True, align=PP_ALIGN.CENTER)
    _multi(slide, l+Inches(0.15), t+Inches(0.55), w-Inches(0.3), h-Inches(0.65),
           body_lines, sz=body_sz, color=C_GRAY_D, align=PP_ALIGN.CENTER, spacing=Pt(3))

def set_ph(slide, idx, text):
    """Set text in a placeholder by index."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return ph
    return None


# ═══════════════════════════════════════════════════════════════════════════
# Slide builders — use template layouts
# ═══════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Cover: use 'Title Slide 1B' (dark bg with title area)."""
    slide = add_slide(prs, 'Title Slide 1B')
    set_ph(slide, 0, "游戏智能 NPC 解决方案")
    set_ph(slide, 1, "基于 Amazon Bedrock & AgentCore 的 AI 驱动动态任务系统")
    set_ph(slide, 10, "Amazon Web Services")
    set_ph(slide, 11, "解决方案详细介绍")
    set_ph(slide, 12, "")

def slide_agenda(prs):
    """Agenda: use 'Agenda Slide 1'."""
    slide = add_slide(prs, 'Agenda Slide 1')
    set_ph(slide, 0, "Agenda")
    body = set_ph(slide, 10, "")
    if body and body.has_text_frame:
        tf = body.text_frame
        items = [
            "01   客户痛点与应用场景 — 传统 NPC 的局限与 AI 智能 NPC 的价值",
            "02   AWS 核心产品介绍 — Amazon Bedrock & AgentCore",
            "03   方案技术架构 — 整体架构、AI 模块架构、行为日志时序",
            "04   Demo 演示 — 智能 NPC 实机演示视频",
            "05   价格分析 — AgentCore & Bedrock 按月成本估算",
            "06   方案总结 — 方案优势对比与正式生产架构",
        ]
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.space_after = Pt(14)

def slide_section(prs, title):
    """Section divider: 'Section Header Option 1' (dark bg)."""
    slide = add_slide(prs, 'Section Header Option 1')
    set_ph(slide, 0, title)

def slide_pain_points(prs):
    """Pain points: use 'Comparison' layout."""
    slide = add_slide(prs, 'Comparison')
    set_ph(slide, 0, "客户痛点")

    # Left header
    ph1 = set_ph(slide, 1, "传统 NPC 实现方式")
    # Left body
    ph2 = set_ph(slide, 2, "")
    if ph2 and ph2.has_text_frame:
        tf = ph2.text_frame
        trad = [
            "NPC 对话基于固定脚本，缺乏变化和沉浸感",
            "任务由策划团队手动录入，开发周期长",
            "玩家行为无法影响 NPC 反应，体验千篇一律",
            "内容更新需要发版，无法动态调整",
            "策划人力成本高，难以规模化生成内容",
        ]
        for i, item in enumerate(trad):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u2716  {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = C_RED
            p.space_after = Pt(8)

    # Right header
    ph3 = set_ph(slide, 3, "AI 智能 NPC 解决方案")
    # Right body
    ph4 = set_ph(slide, 4, "")
    if ph4 and ph4.has_text_frame:
        tf = ph4.text_frame
        ai = [
            "NPC 对话由大语言模型动态生成，自然且个性化",
            "任务根据玩家行为自动生成，无需策划手动设计",
            "基于玩家行为日志分析，NPC 动态调整任务方向",
            "实时更新，无需发版即可优化 NPC 行为",
            "显著降低策划人力成本，提高内容生成效率",
        ]
        for i, item in enumerate(ai):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u2714  {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = C_GREEN
            p.space_after = Pt(8)

def slide_scenario_flow(prs):
    """Core scenario flow — 'Title Only' + custom shapes."""
    slide = add_slide(prs, 'Title Only')
    set_ph(slide, 0, "应用场景 — 核心场景流程")

    def box(l, t, w, h, text, fill, fc=C_WHITE, sz=13):
        _rrect(slide, l, t, w, h, fill=fill)
        _tb(slide, l, t, w, h, text, sz=sz, color=fc, bold=True,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    def diamond(l, t, w, h, text, fill, fc=C_WHITE, sz=12):
        s = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
        _tb(slide, l+Inches(0.1), t+Inches(0.15), w-Inches(0.2), h-Inches(0.3),
            text, sz=sz, color=fc, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Start
    box(Inches(5.2), Inches(1.8), Inches(2.6), Inches(0.55), "玩家进入场景", C_DARK)
    box(Inches(5.2), Inches(2.6), Inches(2.6), Inches(0.55), "键盘操控角色移动", C_BLUE)
    diamond(Inches(5.3), Inches(3.4), Inches(2.4), Inches(1.1), "遭遇对象?", C_GRAY_D)

    _tb(slide, Inches(3.5), Inches(3.6), Inches(1.5), Inches(0.25),
        "\u2190  怪物", sz=11, color=C_RED, bold=True, align=PP_ALIGN.RIGHT)
    _tb(slide, Inches(8.0), Inches(3.6), Inches(1.5), Inches(0.25),
        "NPC  \u2192", sz=11, color=C_BLUE, bold=True, align=PP_ALIGN.LEFT)

    # Left: Battle
    box(Inches(1.0), Inches(3.6), Inches(2.2), Inches(0.55), "自动进入战斗", C_RED)
    diamond(Inches(1.1), Inches(4.5), Inches(2.0), Inches(0.9), "战斗结果?", C_RED, sz=11)

    _tb(slide, Inches(0.2), Inches(4.7), Inches(0.9), Inches(0.2),
        "胜利\u2193", sz=10, color=C_GREEN, bold=True, align=PP_ALIGN.CENTER)
    box(Inches(0.0), Inches(5.5), Inches(1.8), Inches(0.8), "战斗结算\n经验/金币/掉落", C_GREEN, sz=11)

    _tb(slide, Inches(3.0), Inches(4.7), Inches(0.9), Inches(0.2),
        "失败\u2193", sz=10, color=RGBColor(0x99,0x00,0x00), bold=True, align=PP_ALIGN.CENTER)
    box(Inches(2.1), Inches(5.5), Inches(1.8), Inches(0.8), "回到出生点\n恢复满血",
        RGBColor(0x99,0x00,0x00), sz=11)

    # Right: NPC
    box(Inches(9.5), Inches(3.6), Inches(2.5), Inches(0.55), "触发 NPC 对话", C_BLUE)
    _rrect(slide, Inches(9.2), Inches(4.5), Inches(3.2), Inches(1.2),
           fill=C_AMBER_L, border=C_ORANGE)
    _tb(slide, Inches(9.3), Inches(4.55), Inches(3.0), Inches(0.3),
        "\u2728 AI 核心能力", sz=14, color=C_SQUID, bold=True, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(9.3), Inches(4.9), Inches(3.0), Inches(0.7),
        "NPC Agent 分析玩家行为日志\n动态生成个性化任务", sz=12,
        color=C_GRAY_D, align=PP_ALIGN.CENTER)
    box(Inches(9.5), Inches(5.9), Inches(2.5), Inches(0.55), "玩家接取任务", C_GREEN)


def slide_bedrock(prs):
    """'Title Only' + custom cards for Bedrock."""
    slide = add_slide(prs, 'Title and Subtitle')
    set_ph(slide, 0, "Amazon Bedrock")
    set_ph(slide, 10, "全托管的生成式 AI 服务，提供来自领先 AI 公司的高性能基础模型")

    cards = [
        ("多模型选择", ["提供 Claude、Llama 等", "多家顶级模型", "按需选择最佳模型"]),
        ("Converse API", ["统一对话接口", "简化多模型集成", "支持 Tool Use"]),
        ("安全合规", ["数据不用于模型训练", "VPC 私有端点", "SOC / ISO 合规"]),
        ("按量计费", ["按 Token 计费", "无最低消费", "成本可预测"]),
    ]
    for i, (t, b) in enumerate(cards):
        x = Inches(0.5 + i * 3.1)
        _card(slide, x, Inches(2.2), Inches(2.8), Inches(2.0),
              t, b, title_color=C_BLUE, bg=C_BLUE_L, border=C_BLUE)

    # Model info
    _rrect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3),
           fill=C_AMBER_L, border=C_ORANGE)
    _tb(slide, Inches(0.8), Inches(4.6), Inches(6), Inches(0.35),
        "本方案使用的模型", sz=17, color=C_SQUID, bold=True)
    _hline(slide, Inches(0.8), Inches(5.0), Inches(1.0), C_ORANGE, Pt(2))
    _tb(slide, Inches(0.8), Inches(5.15), Inches(11), Inches(0.3),
        "Claude 3.5 Haiku  (us.anthropic.claude-3-5-haiku-20241022-v1:0)",
        sz=15, color=C_DARK, bold=True)
    _multi(slide, Inches(0.8), Inches(5.55), Inches(11), Inches(1.1),
           ["Anthropic 最新轻量级模型，推理速度极快、成本极低",
            "非常适合游戏场景中的实时 NPC 对话和任务生成",
            "单次调用延迟约 3.5 秒，满足游戏交互的实时性要求",
            "200K Token 上下文窗口，可注入丰富的游戏数据作为 context"],
           sz=13, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(5))


def slide_agentcore(prs):
    slide = add_slide(prs, 'Title and Subtitle')
    set_ph(slide, 0, "Amazon Bedrock AgentCore")
    set_ph(slide, 10, "全托管的 AI Agent 部署与运行服务，支持自定义 Agent 容器化部署")

    cards = [
        ("容器化部署", ["支持自定义 Docker 镜像", "灵活部署任意框架", "ARM64 架构支持"]),
        ("VPC 私有模式", ["Agent 运行在客户 VPC", "数据不出网络边界", "私有子网隔离"]),
        ("弹性扩缩", ["按请求量自动扩缩", "无需管理基础设施", "高可用架构"]),
        ("Runtime Endpoint", ["标准化 API 端点", "简化调用集成", "内置监控日志"]),
    ]
    for i, (t, b) in enumerate(cards):
        x = Inches(0.5 + i * 3.1)
        _card(slide, x, Inches(2.2), Inches(2.8), Inches(2.0),
              t, b, title_color=C_BLUE, bg=C_BLUE_L, border=C_BLUE)

    _rrect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3),
           fill=C_AMBER_L, border=C_ORANGE)
    _tb(slide, Inches(0.8), Inches(4.6), Inches(6), Inches(0.35),
        "本方案中的使用方式", sz=17, color=C_SQUID, bold=True)
    _hline(slide, Inches(0.8), Inches(5.0), Inches(1.0), C_ORANGE, Pt(2))
    _multi(slide, Inches(0.8), Inches(5.15), Inches(11), Inches(1.5),
           ["NPC Agent 以 Docker 容器形式部署到 AgentCore Runtime (ARM64)",
            "Agent 运行在 VPC 私有子网，通过 NAT Gateway 访问 Bedrock 和 DynamoDB",
            "游戏服务器通过 InvokeAgentRuntime API 调用 NPC Agent",
            "Agent 预获取 6 张 DynamoDB 表数据，单次调用 Bedrock Converse API 完成推理"],
           sz=13, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(6))


def slide_arch_highlevel(prs):
    slide = add_slide(prs, 'Title and Subtitle')
    set_ph(slide, 0, "整体架构 (High Level)")
    set_ph(slide, 10, "详细架构图请参见 architecture.drawio — Page 1: 当前 Demo 架构")

    # Dev terminal
    _rrect(slide, Inches(0.3), Inches(3.5), Inches(1.5), Inches(0.7), fill=C_DARK)
    _tb(slide, Inches(0.3), Inches(3.5), Inches(1.5), Inches(0.7),
        "开发者终端", sz=12, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, Inches(1.8), Inches(3.55), Inches(1.2), Inches(0.3),
        "SSM 端口转发 \u2192", sz=10, color=C_GRAY_D, align=PP_ALIGN.CENTER)

    # VPC
    _rrect(slide, Inches(3.0), Inches(2.0), Inches(4.5), Inches(4.5), border=C_BLUE, bw=Pt(2))
    _tb(slide, Inches(3.1), Inches(2.05), Inches(3), Inches(0.25),
        "AWS VPC — 私有子网", sz=11, color=C_BLUE, bold=True)
    # EC2
    _rrect(slide, Inches(3.2), Inches(2.5), Inches(4.1), Inches(3.7), fill=C_BLUE_L, border=C_BLUE)
    _tb(slide, Inches(3.3), Inches(2.55), Inches(3), Inches(0.25),
        "Amazon EC2 (Node.js + Express)", sz=11, color=C_BLUE, bold=True)
    _multi(slide, Inches(3.4), Inches(3.0), Inches(3.8), Inches(2.8),
           ["Express.static — 前端静态资源", "REST API / WebSocket",
            "BattleSystem — 战斗系统", "TaskManager — 任务奖励结算",
            "EventEmitter — 事件发射器"], sz=12, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(6))

    # AI Module
    _rrect(slide, Inches(8.0), Inches(2.0), Inches(4.5), Inches(2.8), fill=C_AMBER_L, border=C_ORANGE, bw=Pt(2))
    _tb(slide, Inches(8.1), Inches(2.05), Inches(3), Inches(0.25),
        "AI 模块", sz=11, color=C_SQUID, bold=True)
    _multi(slide, Inches(8.2), Inches(2.5), Inches(4.0), Inches(2.0),
           ["AgentCore Runtime (VPC 私有模式)", "NPC Agent (Docker 容器)",
            "Bedrock Converse API", "Claude 3.5 Haiku"],
           sz=12, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(6))
    _tb(slide, Inches(7.1), Inches(2.8), Inches(1.1), Inches(0.3),
        "NPC 对话请求 \u2192", sz=9, color=C_SQUID, bold=True, align=PP_ALIGN.CENTER)

    # DynamoDB
    _rrect(slide, Inches(8.0), Inches(5.1), Inches(4.5), Inches(1.4), fill=C_GREEN_L, border=C_GREEN, bw=Pt(2))
    _tb(slide, Inches(8.1), Inches(5.15), Inches(3), Inches(0.25),
        "Amazon DynamoDB", sz=11, color=C_GREEN, bold=True)
    _multi(slide, Inches(8.2), Inches(5.5), Inches(4.0), Inches(0.8),
           ["Players | Monsters | NPCs | Items", "Tasks | PlayerEventSummary"],
           sz=12, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(4))
    _tb(slide, Inches(7.1), Inches(5.5), Inches(1.1), Inches(0.3),
        "\u2190 读写 \u2192", sz=9, color=C_GREEN, bold=True, align=PP_ALIGN.CENTER)


def slide_arch_ai_module(prs):
    slide = add_slide(prs, 'Title and Subtitle')
    set_ph(slide, 0, "AI 模块架构 — NPC Agent 内部流程")
    set_ph(slide, 10, "详细架构图请参见 architecture.drawio — Page 3: NPC Agent 内部架构")

    # Pipeline steps
    steps = [
        ("1. 数据预获取", "并行读取 6 张\nDynamoDB 表", C_BLUE),
        ("2. Prompt 组装", "精简数据 +\nSystem Prompt", C_GRAY_D),
        ("3. LLM 调用", "单次 Bedrock\nConverse API", C_ORANGE),
        ("4. JSON 解析", "解析 dialogue\n+ task", C_GRAY_D),
        ("5. 任务校验", "validate_task\n6 项校验规则", C_RED),
        ("6. 写入 DB", "校验通过后\n写入 Tasks 表", C_GREEN),
    ]
    for i, (t, d, c) in enumerate(steps):
        x = Inches(0.25 + i * 2.15)
        _rrect(slide, x, Inches(2.0), Inches(1.9), Inches(1.6), fill=C_WHITE, border=c, bw=Pt(2))
        _rect(slide, x, Inches(2.0), Inches(1.9), Inches(0.06), fill=c)
        _tb(slide, x+Inches(0.1), Inches(2.1), Inches(1.7), Inches(0.35),
            t, sz=13, color=c, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, x+Inches(0.1), Inches(2.5), Inches(1.7), Inches(0.9),
            d, sz=11, color=C_GRAY_D, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            _tb(slide, x+Inches(1.85), Inches(2.5), Inches(0.4), Inches(0.4),
                "\u2192", sz=20, color=C_ORANGE, bold=True, align=PP_ALIGN.CENTER)

    # Data sources
    _tb(slide, Inches(0.5), Inches(3.9), Inches(10), Inches(0.3),
        "预获取数据源 (6 张 DynamoDB 表)", sz=15, color=C_DARK, bold=True)
    _hline(slide, Inches(0.5), Inches(4.22), Inches(1.0), C_GREEN, Pt(2))
    tables = [("Players","玩家状态"), ("PlayerEvent\nSummary","最近 20 条\n行为日志"),
              ("Monsters","怪物字典"), ("Items","道具字典"), ("NPCs","NPC 字典"), ("Tasks","已有任务")]
    for i, (n, d) in enumerate(tables):
        x = Inches(0.25 + i * 2.15)
        _rrect(slide, x, Inches(4.4), Inches(1.9), Inches(0.95), fill=C_GREEN_L, border=C_GREEN)
        _tb(slide, x+Inches(0.05), Inches(4.45), Inches(1.8), Inches(0.35),
            n, sz=11, color=C_GREEN, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, x+Inches(0.05), Inches(4.8), Inches(1.8), Inches(0.4),
            d, sz=10, color=C_GRAY_D, align=PP_ALIGN.CENTER)

    # Validation + Performance
    _tb(slide, Inches(0.5), Inches(5.6), Inches(5), Inches(0.3),
        "任务校验规则 (validate_task)", sz=15, color=C_DARK, bold=True)
    _hline(slide, Inches(0.5), Inches(5.92), Inches(1.0), C_RED, Pt(2))
    rules = [
        "JSON 结构完整性 (title / description / conditions / awards)",
        "条件类型合法性 (kill_monster / collect_item / talk_to_npc / use_item)",
        "target_id 在字典表中存在 | 奖励 item_id 在字典表中存在",
        "数值范围校验 (金币 1-1000, 经验 1-500) | 任务去重",
    ]
    _multi(slide, Inches(0.5), Inches(6.0), Inches(6.0), Inches(1.2),
           rules, sz=11, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(2))

    _rrect(slide, Inches(7.0), Inches(5.6), Inches(5.5), Inches(1.5),
           fill=C_AMBER_L, border=C_ORANGE)
    _tb(slide, Inches(7.2), Inches(5.65), Inches(3), Inches(0.25),
        "性能特征", sz=14, color=C_SQUID, bold=True)
    _hline(slide, Inches(7.2), Inches(5.92), Inches(0.8), C_ORANGE, Pt(2))
    _multi(slide, Inches(7.2), Inches(6.0), Inches(5.0), Inches(1.0),
           ["数据预获取: ~200ms", "LLM 推理 (Claude 3.5 Haiku): ~3.5s",
            "校验 + 写入: ~100ms", "总耗时: ~4 秒"],
           sz=12, color=C_GRAY_D, bullet="\u2022  ", spacing=Pt(3))


def slide_sequence(prs):
    slide = add_slide(prs, 'Title Only')
    set_ph(slide, 0, "行为日志采集与供给时序")

    parts = [
        ("玩家", Inches(0.3), C_DARK),
        ("游戏服务", Inches(2.3), C_BLUE),
        ("Kinesis", Inches(4.5), C_ORANGE),
        ("S3", Inches(6.3), C_GREEN),
        ("DynamoDB", Inches(8.3), C_BLUE),
        ("NPC Agent", Inches(10.5), C_SQUID),
    ]
    for name, x, color in parts:
        _rrect(slide, x, Inches(1.8), Inches(1.6), Inches(0.5), fill=color)
        _tb(slide, x, Inches(1.8), Inches(1.6), Inches(0.5),
            name, sz=12, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _rect(slide, x+Inches(0.78), Inches(2.3), Pt(1), Inches(4.5),
              fill=RGBColor(0xBB,0xBB,0xBB))

    msgs = [
        (0.3, 2.3, "1. 玩家行为 (战斗/道具/...)", 2.6, C_DARK),
        (2.3, 4.5, "2. 发送完整事件 JSON", 3.1, C_ORANGE),
        (4.5, 6.3, "3. Firehose 持久化全量日志", 3.5, C_GREEN),
        (2.3, 8.3, "4. 写入最近事件摘要 (保留50条)", 4.0, C_BLUE),
        (0.3, 2.3, "5. 接近 NPC, 触发对话", 4.6, C_DARK),
        (2.3, 10.5, "6. InvokeAgentRuntime (player_id, npc_id)", 5.1, C_SQUID),
        (10.5, 8.3, "7. get_player_events(limit=20)", 5.6, C_BLUE),
        (8.3, 10.5, "8. 返回最近 20 条行为日志", 6.1, C_SQUID),
    ]
    for frm, to, text, y, color in msgs:
        mn, mx = min(frm, to), max(frm, to)
        _tb(slide, Inches(mn+0.8), Inches(y-0.18), Inches(mx-mn), Inches(0.22),
            text, sz=9, color=color, align=PP_ALIGN.CENTER)
        _rect(slide, Inches(mn+0.8), Inches(y+0.04), Inches(mx-mn), Pt(1.5), fill=color)


def slide_demo(prs):
    slide = add_slide(prs, 'Video or Demo Divider')
    set_ph(slide, 0, "Demo 演示")
    _tb(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.5),
        "[ 请在此处插入 Demo 视频 ]", sz=16, color=C_GRAY_D, align=PP_ALIGN.CENTER)


def slide_pricing(prs):
    slide = add_slide(prs, 'Title Only')
    set_ph(slide, 0, "价格分析 — 月度成本估算")

    _tb(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.3),
        "以下估算基于 Demo 使用场景 (约 1000 次 NPC 对话/月)", sz=12, color=C_GRAY_D)

    # AgentCore card
    _rrect(slide, Inches(0.5), Inches(2.1), Inches(5.8), Inches(3.5), fill=C_WHITE, border=C_ORANGE, bw=Pt(2))
    _rect(slide, Inches(0.5), Inches(2.1), Inches(5.8), Inches(0.5), fill=C_ORANGE)
    _tb(slide, Inches(0.6), Inches(2.12), Inches(5.5), Inches(0.46),
        "Amazon Bedrock AgentCore", sz=16, color=C_WHITE, bold=True)

    ac = ["计费模型: 按 Agent 运行时间计费", "",
          "Runtime 实例 (arm64, 0.5 vCPU + 1GB):", "    ~$0.01 / 小时",
          "    月度 (24x7): ~$7.20", "",
          "Endpoint 请求费: 可忽略不计", "",
          "AgentCore 月度估算: ~$7 - $15"]
    y = Inches(2.75)
    for line in ac:
        if not line: y += Inches(0.1); continue
        is_total = line.startswith("AgentCore 月度")
        _tb(slide, Inches(0.8), y, Inches(5.2), Inches(0.22),
            line, sz=11, color=C_SQUID if is_total else C_GRAY_D, bold=is_total)
        y += Inches(0.22)

    # Bedrock card
    _rrect(slide, Inches(6.8), Inches(2.1), Inches(5.8), Inches(3.5), fill=C_WHITE, border=C_BLUE, bw=Pt(2))
    _rect(slide, Inches(6.8), Inches(2.1), Inches(5.8), Inches(0.5), fill=C_BLUE)
    _tb(slide, Inches(6.9), Inches(2.12), Inches(5.5), Inches(0.46),
        "Amazon Bedrock (Claude 3.5 Haiku)", sz=16, color=C_WHITE, bold=True)

    br = ["计费模型: 按输入/输出 Token 计费", "",
          "Claude 3.5 Haiku 定价:", "    输入: $0.80 / 百万 Token",
          "    输出: $4.00 / 百万 Token", "",
          "Demo 场景 (1000 次对话/月):", "    月输入 ~200 万 Token = $1.60",
          "    月输出 ~50 万 Token = $2.00", "",
          "Bedrock 月度估算: ~$3.60"]
    y = Inches(2.75)
    for line in br:
        if not line: y += Inches(0.1); continue
        is_total = line.startswith("Bedrock 月度")
        _tb(slide, Inches(7.1), y, Inches(5.2), Inches(0.22),
            line, sz=11, color=C_BLUE if is_total else C_GRAY_D, bold=is_total)
        y += Inches(0.22)

    # Total
    _rrect(slide, Inches(0.5), Inches(5.8), Inches(12.1), Inches(1.1),
           fill=C_AMBER_L, border=C_ORANGE, bw=Pt(2))
    _tb(slide, Inches(0.8), Inches(5.85), Inches(5), Inches(0.3),
        "月度总成本估算 (Demo 场景)", sz=16, color=C_SQUID, bold=True)
    _tb(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.25),
        "AgentCore: ~$7-15  |  Bedrock: ~$3.60  |  DynamoDB (按需): ~$1-5  |  其他 (EC2, VPC): 按现有基础设施",
        sz=11, color=C_GRAY_D)
    _tb(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.3),
        "AI 模块月度总成本 (AgentCore + Bedrock): 约 $10 - $20 — 成本极低",
        sz=14, color=C_SQUID, bold=True)


def slide_comparison(prs):
    slide = add_slide(prs, 'Title Only')
    set_ph(slide, 0, "方案总结 — 优势对比")

    cols = [(Inches(0.5), Inches(2.0)), (Inches(2.5), Inches(4.2)), (Inches(6.7), Inches(5.6))]
    for (x, w), t in zip(cols, ["对比维度", "传统方式", "AI 智能 NPC 方案"]):
        _rect(slide, x, Inches(1.8), w, Inches(0.48), fill=C_DARK)
        _tb(slide, x, Inches(1.8), w, Inches(0.48),
            t, sz=13, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    rows = [
        ("NPC 对话", "固定脚本，千篇一律", "LLM 动态生成，自然且个性化"),
        ("任务生成", "策划手动设计和录入", "AI 根据玩家行为自动生成"),
        ("内容更新", "需要发版更新", "实时动态调整，无需发版"),
        ("玩家体验", "缺乏沉浸感", "个性化体验，NPC 理解玩家处境"),
        ("开发成本", "大量策划人力，周期长", "一次开发，AI 持续生成内容"),
        ("可扩展性", "内容量与人力成正比", "AI 无限生成，边际成本趋零"),
        ("数据驱动", "无法利用玩家行为数据", "行为日志驱动决策，持续优化"),
        ("运营成本", "高 (人力为主)", "极低 (~$10-20/月 AI 服务费)"),
    ]
    for i, (dim, old, new) in enumerate(rows):
        y = Inches(2.28 + i * 0.52)
        bg = C_GRAY_L if i % 2 == 0 else C_WHITE
        for x, w in [(Inches(0.5), Inches(2.0)), (Inches(2.5), Inches(4.2)), (Inches(6.7), Inches(5.6))]:
            _rect(slide, x, y, w, Inches(0.49), fill=bg)
        _tb(slide, Inches(0.6), y+Inches(0.08), Inches(1.8), Inches(0.35),
            dim, sz=12, color=C_DARK, bold=True)
        _tb(slide, Inches(2.6), y+Inches(0.08), Inches(4.0), Inches(0.35),
            old, sz=11, color=C_RED)
        _tb(slide, Inches(6.8), y+Inches(0.08), Inches(5.4), Inches(0.35),
            new, sz=11, color=C_GREEN)


def slide_prod_arch(prs):
    slide = add_slide(prs, 'Title and Subtitle')
    set_ph(slide, 0, "正式生产架构（目标）")
    set_ph(slide, 10, "详细架构图请参见 architecture.drawio — Page 2: 正式生产架构")

    # User -> CloudFront
    _rrect(slide, Inches(0.3), Inches(2.5), Inches(1.3), Inches(0.6), fill=C_DARK)
    _tb(slide, Inches(0.3), Inches(2.5), Inches(1.3), Inches(0.6),
        "用户", sz=11, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, Inches(1.6), Inches(2.55), Inches(0.5), Inches(0.3),
        "\u2192", sz=16, color=C_GRAY_D, align=PP_ALIGN.CENTER)

    _rrect(slide, Inches(2.0), Inches(2.5), Inches(1.5), Inches(0.6), fill=C_PURPLE)
    _tb(slide, Inches(2.0), Inches(2.5), Inches(1.5), Inches(0.6),
        "CloudFront", sz=11, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, Inches(2.2), Inches(3.15), Inches(1.1), Inches(0.2),
        "\u2193 静态资源", sz=9, color=C_GRAY_D, align=PP_ALIGN.CENTER)
    _rrect(slide, Inches(2.0), Inches(3.4), Inches(1.5), Inches(0.5), fill=C_GREEN)
    _tb(slide, Inches(2.0), Inches(3.4), Inches(1.5), Inches(0.5),
        "S3 前端", sz=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, Inches(3.5), Inches(2.55), Inches(0.8), Inches(0.3),
        "/api /ws \u2192", sz=9, color=C_GRAY_D, align=PP_ALIGN.CENTER)

    # VPC
    _rrect(slide, Inches(4.2), Inches(2.0), Inches(4.3), Inches(4.2), border=C_BLUE, bw=Pt(2))
    _tb(slide, Inches(4.3), Inches(2.05), Inches(2), Inches(0.2), "VPC", sz=10, color=C_BLUE, bold=True)
    _rrect(slide, Inches(4.4), Inches(2.4), Inches(1.5), Inches(0.55), fill=C_PURPLE)
    _tb(slide, Inches(4.4), Inches(2.4), Inches(1.5), Inches(0.55),
        "ALB", sz=11, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _rrect(slide, Inches(4.4), Inches(3.2), Inches(3.8), Inches(2.7), fill=C_BLUE_L, border=C_BLUE)
    _tb(slide, Inches(4.5), Inches(3.25), Inches(2), Inches(0.2),
        "EC2 (私有子网)", sz=10, color=C_BLUE, bold=True)
    _multi(slide, Inches(4.6), Inches(3.55), Inches(3.4), Inches(2.0),
           ["REST API / WebSocket", "战斗系统", "任务管理", "事件发射器"],
           sz=10, color=C_GRAY_D, bullet="\u2022 ", spacing=Pt(3))

    # AI Module
    _rrect(slide, Inches(9.0), Inches(2.0), Inches(3.5), Inches(2.0), fill=C_AMBER_L, border=C_ORANGE, bw=Pt(2))
    _tb(slide, Inches(9.1), Inches(2.05), Inches(2), Inches(0.2), "AI 模块", sz=10, color=C_SQUID, bold=True)
    _multi(slide, Inches(9.2), Inches(2.4), Inches(3.2), Inches(1.4),
           ["AgentCore Runtime", "NPC Agent (Docker)", "Bedrock Claude 3.5 Haiku"],
           sz=10, color=C_GRAY_D, bullet="\u2022 ", spacing=Pt(3))
    _tb(slide, Inches(8.1), Inches(3.0), Inches(1.0), Inches(0.3),
        "NPC\u2192", sz=9, color=C_SQUID, align=PP_ALIGN.CENTER)

    # Pipeline
    _rrect(slide, Inches(9.0), Inches(4.3), Inches(3.5), Inches(1.3), fill=C_AMBER_L, border=C_ORANGE, bw=Pt(2))
    _tb(slide, Inches(9.1), Inches(4.35), Inches(2), Inches(0.2), "数据流管道", sz=10, color=C_SQUID, bold=True)
    _multi(slide, Inches(9.2), Inches(4.65), Inches(3.2), Inches(0.8),
           ["Kinesis Data Stream", "Kinesis Firehose \u2192 S3"],
           sz=10, color=C_GRAY_D, bullet="\u2022 ", spacing=Pt(3))
    _tb(slide, Inches(8.1), Inches(4.7), Inches(1.0), Inches(0.3),
        "事件\u2192", sz=9, color=C_SQUID, align=PP_ALIGN.CENTER)

    # DynamoDB + S3
    _rrect(slide, Inches(4.2), Inches(6.4), Inches(3.5), Inches(0.6), fill=C_GREEN_L, border=C_GREEN)
    _tb(slide, Inches(4.2), Inches(6.4), Inches(3.5), Inches(0.6),
        "DynamoDB (6 tables)", sz=12, color=C_GREEN, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _rrect(slide, Inches(9.0), Inches(6.4), Inches(3.5), Inches(0.6), fill=C_GREEN_L, border=C_GREEN)
    _tb(slide, Inches(9.0), Inches(6.4), Inches(3.5), Inches(0.6),
        "S3 行为日志归档", sz=12, color=C_GREEN, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _tb(slide, Inches(5.5), Inches(6.05), Inches(0.5), Inches(0.3), "\u2193", sz=14, color=C_GREEN, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(10.3), Inches(5.95), Inches(0.5), Inches(0.3), "\u2193", sz=14, color=C_GREEN, align=PP_ALIGN.CENTER)


def slide_thankyou(prs):
    slide = add_slide(prs, 'Thank You Option 1')
    # Placeholders 10-15 are name/title pairs
    set_ph(slide, 10, "")
    set_ph(slide, 11, "")
    set_ph(slide, 12, "")
    set_ph(slide, 13, "")
    set_ph(slide, 14, "")
    set_ph(slide, 15, "")


# ═══════════════════════════════════════════════════════════════════════════
def main():
    print(f"Loading template: {TEMPLATE_PATH}")
    prs = Presentation(TEMPLATE_PATH)
    print(f"  Slides before: {len(prs.slides)}")

    delete_all_slides(prs)
    print(f"  Slides after clear: {len(prs.slides)}")

    # Build all slides
    slide_title(prs);          print("  [1]  Title")
    slide_agenda(prs);         print("  [2]  Agenda")
    slide_section(prs, "01  客户痛点与应用场景"); print("  [3]  Section 1")
    slide_pain_points(prs);    print("  [4]  Pain points")
    slide_scenario_flow(prs);  print("  [5]  Scenario flow")
    slide_section(prs, "02  AWS 核心产品介绍");   print("  [6]  Section 2")
    slide_bedrock(prs);        print("  [7]  Bedrock")
    slide_agentcore(prs);      print("  [8]  AgentCore")
    slide_section(prs, "03  方案技术架构");       print("  [9]  Section 3")
    slide_arch_highlevel(prs); print("  [10] High-level arch")
    slide_arch_ai_module(prs); print("  [11] AI module arch")
    slide_sequence(prs);       print("  [12] Sequence diagram")
    slide_section(prs, "04  Demo 演示");          print("  [13] Section 4")
    slide_demo(prs);           print("  [14] Demo placeholder")
    slide_section(prs, "05  价格分析");           print("  [15] Section 5")
    slide_pricing(prs);        print("  [16] Pricing")
    slide_section(prs, "06  方案总结");           print("  [17] Section 6")
    slide_comparison(prs);     print("  [18] Comparison table")
    slide_prod_arch(prs);      print("  [19] Production arch")
    slide_thankyou(prs);       print("  [20] Thank You")

    prs.save(OUTPUT_PATH)
    print(f"\nSaved: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
